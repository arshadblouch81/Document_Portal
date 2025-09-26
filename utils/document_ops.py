from __future__ import annotations

from pathlib import Path
import textwrap
from typing import Iterable, List

from langchain.schema import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
from logger.custom_logger import CustomLogger
from exception.custom_exception import DocumentPortalException
import sqlite3
from langchain_community.document_loaders  import (
    PyPDFLoader,
    Docx2txtLoader,
    TextLoader,
    UnstructuredMarkdownLoader,
    UnstructuredPowerPointLoader,
    UnstructuredExcelLoader,
    CSVLoader,
    UnstructuredFileLoader
)
from PIL import Image
import pytesseract
import pdfplumber
import fitz  # PyMuPDF
from PIL import Image

import io
from docx import Document as DocxDocument
import docx2txt
import pandas as pd
from pptx import Presentation
from utils.file_io import  encode_image_to_base64,preprocess_image
from utils.model_loader import ModelLoader
from langchain_core.messages import HumanMessage

log = CustomLogger().get_logger(__name__)


SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".txt", ".md", ".ppt", ".pptx", ".xlsx", ".csv", ".sql", ".json", ".jpg", ".png", ".jpeg", ".gif", ".tiff", ".bmp", ".webp", ".svg"}

from langchain_community.document_loaders import JSONLoader

def load_documents(paths: Iterable[Path]) -> List[Document]:
    """Load docs using appropriate loader based on extension."""
    docs: List[Document] = []
    table_docs: List[Document] = []
    fig_docs: List[Document] = []

    try:
        for p in paths:
            ext = p.suffix.lower()
            loader = None

            if ext == ".pdf":
                loader = PyPDFLoader(str(p))
            elif ext == ".docx":
                loader = Docx2txtLoader(str(p))
            elif ext == ".txt":
                loader = TextLoader(str(p), encoding="utf-8")
            elif ext == ".md":
                loader = UnstructuredMarkdownLoader(str(p))
            elif ext == ".ppt" or ext == ".pptx":
                loader = UnstructuredPowerPointLoader(str(p))
            elif ext == ".xlsx":
                loader = UnstructuredExcelLoader(str(p))
            elif ext == ".csv":
                loader = CSVLoader(file_path=str(p))
            elif ext == ".sql":
                docs_list = load_sql_file(str(p))
                docs.extend(docs_list)
                continue
            elif ext == ".json":
                # Customize jq_schema based on your JSON structure
                loader = JSONLoader(
                    file_path=str(p),
                    jq_schema=".[] | .content",  # Example: extract 'content' from each message
                    text_content=False  # Set to True if you want raw text
                )
            elif ext in [".png", ".jpg", ".jpeg", ".bmp", ".tiff"]:
                text = read_image_file(p)
                docs.append(Document(page_content=text, metadata={"source": str(p)}))
                continue
            else:
                log.warning("Unsupported extension skipped", path=str(p))
                continue

            if loader:
                docs.extend(loader.load())
            if table_docs:
                docs.extend(table_docs)
            if fig_docs:
                docs.extend(fig_docs)

        log.info("Documents loaded", count=len(docs))
        return docs

    except Exception as e:
        log.error("Failed loading documents", error=str(e))
        raise DocumentPortalException("Error loading documents", e) from e
    
def load_sql_file(file_path: str) -> str:
    conn = sqlite3.connect('my_database.db')
       
    cursor = conn.cursor()   
    cursor.execute("""
        SELECT name
        FROM sqlite_master 
        WHERE type='table' AND name NOT LIKE 'sqlite_%'
    """)
    tables = cursor.fetchall()  # Get all table names
    
    documents: list[Document] = []

    for table in tables:
        table_name = table[0]
        try:
            cursor.execute(f"SELECT * FROM {table_name}")
            rows = cursor.fetchall()
            columns = [desc[0] for desc in cursor.description]
            for idx, row in enumerate(rows[:10], start=1):
                # Format each column-value pair with indentation for clarity
                content_lines = [f"{col}: {val}" for col, val in zip(columns, row)]
                content = "\n".join(content_lines)

                # Create Document with metadata including row index
                doc = Document(
                    page_content=content,
                    metadata={
                        "source_table": table_name,
                        "row_index": idx
                    }
                )
                documents.append(doc)

            # for row in rows:
            #     content = "\n".join(f"{col}: {val}" for col, val in zip(columns, row))
            #     doc = Document(
            #         page_content=content,
            #         metadata={"source_table": table_name}
            #     )
            #     documents.append(doc)

        except Exception as e:
            print(f"Failed to fetch data from table '{table_name}': {e}")

    cursor.close()    
    conn.close()
    return documents


   
def concat_for_analysis(docs: List[Document]) -> str:
    parts = []
    for d in docs:
        src = d.metadata.get("source") or d.metadata.get("file_path") or "unknown"
        parts.append(f"\n--- SOURCE: {src} ---\n{d.page_content}")
    return "\n".join(parts)

def concat_for_comparison(ref_docs: List[Document], act_docs: List[Document]) -> str:
    left = concat_for_analysis(ref_docs)
    right = concat_for_analysis(act_docs)
    return f"<<REFERENCE_DOCUMENTS>>\n{left}\n\n<<ACTUAL_DOCUMENTS>>\n{right}"

#loading Documents from various table and images in

def extract_tables_from_pdf(path):
    tables = []
    with pdfplumber.open(path) as pdf:
        for page in pdf.pages:
            page_tables = page.extract_tables()
            for table in page_tables:
                tables.append(table)
    return tables

def extract_images_from_pdf(path):
    doc = fitz.open(path)
    texts = []
    for page in doc:
        for img in page.get_images(full=True):
            base_image = doc.extract_image(img[0])
            image = Image.open(io.BytesIO(base_image["image"]))
            text = pytesseract.image_to_string(image)
            texts.append(text)
    return texts

def extract_tables_from_docx(path):
    doc = DocxDocument(path)
    tables = []
    for table in doc.tables:
        rows = []
        for row in table.rows:
            rows.append([cell.text.strip() for cell in row.cells])
        tables.append(rows)
    return tables

def extract_images_from_docx(path):
    temp_dir = "temp_images"
    docx2txt.process(path, temp_dir)
    texts = []
    for img_path in Path(temp_dir).glob("*.*"):
        image = Image.open(img_path)
        text = pytesseract.image_to_string(image)
        texts.append(text)
    return texts

def extract_tables_from_excel(path):
    xls = pd.ExcelFile(path)
    tables = []
    for sheet_name in xls.sheet_names:
        df = xls.parse(sheet_name)
        tables.append(df.to_dict(orient="records"))
    return tables

def extract_tables_from_pptx(path):
    prs = Presentation(path)
    tables = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                rows = []
                for row in table.rows:
                    rows.append([cell.text.strip() for cell in row.cells])
                tables.append(rows)
    return tables
def extract_images_from_pptx(path):
    prs = Presentation(path)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == 13:  # Picture
                image = shape.image
                image_bytes = image.blob
                img = Image.open(io.BytesIO(image_bytes))
                text = pytesseract.image_to_string(img)
                texts.append(text)
    return texts

def read_image_file( file_path: str) -> str:
        try:
            # Load image
            
            pil_image = preprocess_image(file_path)
            image_base64 = encode_image_to_base64(image=pil_image, format=pil_image.format or "PNG")

            # Optional: Embed as data URI
            # data_uri = f"data:image/png;base64,{image_base64}"

            # Construct multimodal message
            mesage = HumanMessage(content=[
                {
                    "type": "text",
                    "text": "Extract all visible text from this image. Include handwritten and printed text."
                },
                {
                    "type": "image_url",
                    "image_url": {
                        "url": f"data:image/png;base64,{image_base64}"
                    }
                }
            ])
            loader = ModelLoader()
            model = loader.load_img_reader_llm()
            response = model.invoke([mesage])
            raw_text = response.content

            # # Extract text using OCR
            # model = self.model_loader.load_img_reader_llm()
            # raw_text = model.generate_content([
            #             image,
            #             "Extract all visible text from this image. Include handwritten and printed text."
            #         ])

            # raw_text = pytesseract.image_to_string(image)

            # Split text into chunks of ~500 characters (adjust as needed)
            wrapped_chunks = textwrap.wrap(raw_text, width=500)
            text_chunks = []

            for i, chunk in enumerate(wrapped_chunks):
                text_chunks.append(f"\n--- Page {i + 1} ---\n{chunk}")

            text = "\n".join(text_chunks)

            log.info("Image read successfully", file_path=file_path, pages=len(text_chunks))
            return text

        except Exception as e:
            log.error("Failed to read image", error=str(e), file_path=file_path)
            raise DocumentPortalException(f"Could not process image: {file_path}", e) from e

def split( docs: List[Document], chunk_size=1000, chunk_overlap=200) -> List[Document]:
    splitter = RecursiveCharacterTextSplitter(chunk_size=chunk_size, chunk_overlap=chunk_overlap)
    chunks = splitter.split_documents(docs)
    log.info("Documents split", chunks=len(chunks), chunk_size=chunk_size, overlap=chunk_overlap)
    return chunks       
     
if __name__ == "__main__":
    paths = [Path("D:\LLMOPS Industry Projects\document_portal\data\email_text.json")]
    #Path("sample.docx"), Path("sample.xlsx"), Path("sample.pptx"), Path("image1.png"
    docs = load_documents(paths)
    chunk_size: int = 1000
    chunk_overlap: int = 200
    k: int = 5
    if not docs:
        raise ValueError("No valid documents loaded")
    
    chunks = split(docs, chunk_size=chunk_size, chunk_overlap=chunk_overlap)
    
    print(chunks[0])
    
    texts = [c.page_content for c in chunks]
    metas = [c.metadata for c in chunks]
    
    
    # try:
    #     vs = fm.load_or_create(texts=texts, metadatas=metas)
    # except Exception:
    #     vs = fm.load_or_create(texts=texts, metadatas=metas)
        
    # added = fm.add_documents(chunks)
    # # log.info("FAISS index updated", added=added, index=str(self.faiss_dir))

    # retriever = vs.as_retriever(search_type="similarity", search_kwargs={"k": k})
